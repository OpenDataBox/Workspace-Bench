"""
validator.py — LLM 后验证模块

负责:
  - validate_content_relevance: 用 LLM 判断下载的文件内容是否与原始文件标题相关
  - extract_text_preview: 从不同格式文件中提取前 N 字作为预览

这是一个独立模块，可单独替换验证策略。
"""

import os
import time
import random
from openai import OpenAI

from pipeline.config import (
    LLM_CONFIG,
    VALIDATOR_LLM_CONFIG, MAX_LLM_RETRIES,
)
from pipeline.utils import log_event


# ================= 后验证 LLM 客户端 =================

_validator_config = VALIDATOR_LLM_CONFIG or LLM_CONFIG

_validator_client = OpenAI(
    api_key=_validator_config["api_key"],
    base_url=_validator_config["base_url"],
)


# ================= 文本预览提取 =================

def extract_text_preview(file_path, max_chars=1500):
    """
    从文件中提取前 max_chars 个字符作为内容预览。

    支持格式: PDF, DOCX, XLSX, PPTX, TXT/DOC(纯文本回退)

    如果对应的库没有安装，回退到读取原始字节中的可读文本。

    Args:
        file_path: 文件路径
        max_chars: 最多提取的字符数

    Returns:
        str: 提取到的文本预览，提取失败返回空字符串
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if ext == '.pdf':
            text = _extract_pdf(file_path, max_chars)
        elif ext in ('.docx',):
            text = _extract_docx(file_path, max_chars)
        elif ext in ('.xlsx', '.xls'):
            text = _extract_xlsx(file_path, max_chars)
        elif ext in ('.pptx',):
            text = _extract_pptx(file_path, max_chars)
        elif ext in ('.txt', '.doc'):
            text = _extract_plain(file_path, max_chars)
    except Exception:
        pass

    # 如果专用提取器失败，尝试通用的字节级文本提取
    if not text.strip():
        text = _extract_raw_text(file_path, max_chars)

    return text[:max_chars].strip()


def _extract_pdf(path, max_chars):
    """尝试用 PyPDF2 或 pdfplumber 提取 PDF 文本。"""
    # 尝试 PyPDF2
    try:
        import PyPDF2
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            texts = []
            for page in reader.pages[:5]:  # 最多前 5 页
                t = page.extract_text() or ""
                texts.append(t)
                if sum(len(x) for x in texts) > max_chars:
                    break
            return "\n".join(texts)
    except ImportError:
        pass

    # 尝试 pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            texts = []
            for page in pdf.pages[:5]:
                t = page.extract_text() or ""
                texts.append(t)
                if sum(len(x) for x in texts) > max_chars:
                    break
            return "\n".join(texts)
    except ImportError:
        pass

    return ""


def _extract_docx(path, max_chars):
    """用 python-docx 提取 DOCX 文本。"""
    try:
        import docx
        doc = docx.Document(path)
        texts = []
        for para in doc.paragraphs:
            texts.append(para.text)
            if sum(len(x) for x in texts) > max_chars:
                break
        return "\n".join(texts)
    except ImportError:
        return ""


def _extract_xlsx(path, max_chars):
    """用 openpyxl 提取 XLSX 前几行文本。"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames[:3]:  # 最多前 3 个 sheet
            ws = wb[sheet]
            texts.append(f"[Sheet: {sheet}]")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                texts.append(" | ".join(cells))
                row_count += 1
                if row_count >= 20 or sum(len(x) for x in texts) > max_chars:
                    break
        wb.close()
        return "\n".join(texts)
    except ImportError:
        return ""


def _extract_pptx(path, max_chars):
    """用 python-pptx 提取 PPTX 幻灯片文本。"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        texts = []
        for i, slide in enumerate(prs.slides[:10]):  # 最多前 10 页
            texts.append(f"[Slide {i+1}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            if sum(len(x) for x in texts) > max_chars:
                break
        return "\n".join(texts)
    except ImportError:
        return ""


def _extract_plain(path, max_chars):
    """读取纯文本文件。"""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(max_chars)
    except Exception:
        return ""


def _extract_raw_text(path, max_chars):
    """
    兜底方案：从文件原始字节中提取可打印的 UTF-8 文本片段。
    不依赖任何第三方库。
    """
    try:
        with open(path, 'rb') as f:
            raw = f.read(max_chars * 3)  # 多读一些字节

        decoded = raw.decode('utf-8', errors='ignore')
        # 只保留可打印字符和换行
        printable = ''.join(
            c for c in decoded
            if c.isprintable() or c in '\n\r\t'
        )
        # 过滤掉过短的片段（可能是二进制噪声）
        lines = [l for l in printable.split('\n') if len(l.strip()) > 3]
        return "\n".join(lines[:50])
    except Exception:
        return ""


# ================= LLM 后验证 =================

def validate_content_relevance(file_path, original_stem, log_file):
    """
    用 LLM 判断下载的文件内容是否与原始文件标题相关。

    流程:
      1. 从下载的文件中提取文本预览
      2. 构造验证 prompt，把文本预览和原始文件名一起发给 LLM
      3. LLM 返回 RELEVANT / IRRELEVANT 判断

    Args:
        file_path: 下载后的文件路径
        original_stem: 原始空文件的文件名（不含后缀），如 "竞品分析报告"
        log_file: 日志文件路径

    Returns:
        (is_relevant: bool, reason: str)
    """
    # Step 1: 提取预览
    preview = extract_text_preview(file_path)

    if not preview or len(preview.strip()) < 20:
        log_event("      [后验] 无法提取有效文本预览，放行。", log_file)
        return True, "无法提取文本，默认放行"

    # Step 2: 构造验证 prompt
    prompt = f"""你是一个文档内容审核员。请判断以下文件内容是否与目标文件标题相关。

目标文件标题: 「{original_stem}」

文件内容预览（前 1500 字）:
---
{preview[:1500]}
---

请判断:
1. 这个文件的实际内容是否与标题「{original_stem}」的语义相关？
2. 这是一份有实质内容的文档，还是一个错误页面/登录页/无关文件？

只返回以下格式（不要有任何其他内容）:
{{"relevant": true, "reason": "简短说明相关原因"}}
或
{{"relevant": false, "reason": "简短说明不相关原因"}}"""

    # Step 3: 调用 LLM 判断
    for attempt in range(MAX_LLM_RETRIES + 1):
        try:
            resp = _validator_client.chat.completions.create(
                model=_validator_config["model"],
                messages=[{"role": "user", "content": prompt}],
                temperature=0.1,  # 极低温度，确保判断稳定
            )
            raw = resp.choices[0].message.content.strip()

            result = _parse_validation_response(raw)
            if result is not None:
                is_relevant, reason = result
                log_event(
                    f"      [后验] {'相关' if is_relevant else '不相关'}: {reason}",
                    log_file,
                )
                return is_relevant, reason
            else:
                log_event(
                    f"      [后验] 响应解析失败 (尝试 {attempt+1}): {raw[:100]}", log_file)

        except Exception as e:
            log_event(
                f"      [后验] 调用异常 (尝试 {attempt+1}): {str(e)[:80]}", log_file)

        if attempt < MAX_LLM_RETRIES:
            time.sleep(random.uniform(1.0, 2.0))

    # 全部失败，默认放行（宁可多留，不错杀）
    log_event("      [后验] 验证失败，默认放行。", log_file)
    return True, "验证失败，默认放行"


def _parse_validation_response(raw_text):
    """
    解析 LLM 的验证响应。

    Returns:
        (is_relevant: bool, reason: str) or None
    """
    import re
    import json

    # 去除 markdown 代码块
    cleaned = raw_text
    cleaned = re.sub(r'^```(?:json)?\s*', '', cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r'```\s*$', '', cleaned, flags=re.MULTILINE)
    cleaned = cleaned.strip()

    # 尝试 JSON
    try:
        obj = json.loads(cleaned)
        if isinstance(obj, dict) and 'relevant' in obj:
            return (bool(obj['relevant']), obj.get('reason', ''))
    except (json.JSONDecodeError, ValueError):
        pass

    # 兜底：关键词匹配
    lower = raw_text.lower()
    if '"relevant": true' in lower or '"relevant":true' in lower:
        return (True, "关键词匹配: relevant=true")
    if '"relevant": false' in lower or '"relevant":false' in lower:
        return (False, "关键词匹配: relevant=false")

    return None
